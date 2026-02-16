#!/usr/bin/env python3
"""
Generate V3 PPTX from rendered slide PNGs.

Defaults are quality-first (master profile):
- Embed source PNG files directly
- Attach speaker notes and alt text
- No transitions

Optional delivery profile:
- Convert source PNG files to JPEG (`--image-format jpeg`)
- Optional max dimension clamp (`--max-dimension`)
"""

from __future__ import annotations

import argparse
import json
import shutil
import subprocess
import tempfile
from datetime import datetime
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.oxml.ns import qn
    from pptx.oxml.xmlchemy import OxmlElement
    from pptx.util import Inches
except ImportError:
    print("Run: pip install python-pptx")
    raise SystemExit(1)

TITLE_FOR = {
    "v3-Slide01TitleV3.png": "Slide 1 — The Guardian Arrives",
    "v3-Slide02ProblemV3.png": "Slide 2 — The Coordination Gap",
    "v3-Slide02ProblemOptionA.png": "Slide 2 — The Coordination Gap",
    "v3-Slide02ProblemOptionB.png": "Slide 2 — The Coordination Gap",
    "v3-Slide02ProblemOptionC.png": "Slide 2 — The Coordination Gap",
    "v3-Slide03WhyNowV3.png": "Slide 3 — Why Now",
    "v3-Slide04SolutionV3.png": "Slide 4 — The Unified AI Backbone",
    "v3-Slide05DifferentiationV3.png": "Slide 5 — Beyond Aggregation",
    "v3-Slide06BusinessV3.png": "Slide 6 — Roadmap & Governance",
    "v3-Slide07DemoV3.png": "Slide 7 — Introduction Video",
    "v3-Slide08SummaryV3.png": "Slide 8 — Summary",
    "v3-Slide09EpilogueV3.png": "Slide 9 — Epilogue",
    "v3-Slide10AppendixV3.png": "Slide 10 — Appendix",
    "v3-Slide11FinModelV3.png": "Slide 11 — Capital-Efficient Monetization",
}


def load_speaker_notes() -> list[str]:
    notes_path = Path(__file__).resolve().parent.parent / "src" / "shared" / "speaker-notes.json"
    if notes_path.exists():
        with open(notes_path, encoding="utf-8") as f:
            notes = json.load(f)
        if isinstance(notes, list) and len(notes) >= 11:
            return notes
        print(f"WARNING: speaker-notes.json has {len(notes)} entries, expected >= 11. Using inline fallback.")
    else:
        print(f"WARNING: {notes_path} not found. Using inline fallback.")
    return [
        "Poseidon.AI is our trusted AI-native money platform. Today I will walk through the four engines and show how they operate as one system.",
        "The hidden coordination tax is real: around $133 per month in user loss, $12.5B annual fraud loss, and major overdraft burden. Mint proved visibility alone is not enough. Users still act as the integrator.",
        "Three forces converge now: open banking standards, AI economics improving, and user expectations shifting to proactive guidance. 2025 is the window for AI-native financial services.",
        "Our architecture is a closed loop: Protect, Grow, Execute, Govern. Each engine has a clear role, and together they deliver auditable decisions, not isolated insights.",
        "We match baseline fintech features, then differentiate on five critical capabilities, especially governance-by-design. That governance moat is difficult to replicate.",
        "Execution is phased across four stages. We track hard gates like 9K users, 70% precision, and 99.9% uptime to prove break-even and operational readiness.",
        "Unit economics justify the business: customer value exceeds acquisition cost by 17x on Plus and 77x on Pro, with blended acquisition cost under $10 driven by organic-first growth. Gross margin is 87% at scale, in line with typical SaaS. We reach operating profitability at Month 15 on $25M total capital, with Series A triggered at $300K monthly revenue in Month 9. By Month 36 we project 2.2 million monthly users, 717K paid, and $96M in annual revenue.",
        "The demo follows a 30-second journey: unified dashboard, Protect with explainability, Execute with human approval, and Govern with audit logging.",
        "The strategy rests on three pillars: governance, architecture, and business model. Economics target about $235 yearly user savings, 77% margin, and roughly 8x value-to-cost.",
        "We close with the team and operating principle: deterministic compute, explainable AI, and explicit human approval before action. The QR code links to the live demo.",
        "Appendix: supporting assumptions, references, and calculation details for diligence review.",
    ]


def clamp_jpeg_quality(raw_quality: int) -> int:
    return max(1, min(100, raw_quality))


def slide02_name(variant: str) -> str:
    if variant == "optionA":
        return "v3-Slide02ProblemOptionA.png"
    if variant == "optionB":
        return "v3-Slide02ProblemOptionB.png"
    if variant == "optionC":
        return "v3-Slide02ProblemOptionC.png"
    return "v3-Slide02ProblemV3.png"


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--slide02",
        choices=["optionA", "optionB", "optionC", "default"],
        default="default",
        help="Select Slide02 variant file.",
    )
    parser.add_argument(
        "--image-format",
        choices=["png", "jpeg"],
        default="png",
        help="Embed images as PNG (master) or JPEG (delivery).",
    )
    parser.add_argument(
        "--jpeg-quality",
        type=int,
        default=82,
        help="JPEG quality (1-100) when --image-format jpeg.",
    )
    parser.add_argument(
        "--max-dimension",
        type=int,
        default=None,
        help="Optional max image dimension (long side) before embedding.",
    )
    parser.add_argument(
        "--output",
        type=Path,
        default=None,
        help="Optional output PPTX path. Default depends on profile.",
    )
    parser.add_argument(
        "--transitions",
        action="store_true",
        help="Enable fade transitions (default: off).",
    )

    notes_group = parser.add_mutually_exclusive_group()
    notes_group.add_argument(
        "--notes",
        dest="notes",
        action="store_true",
        help="Attach speaker notes (default).",
    )
    notes_group.add_argument(
        "--no-notes",
        dest="notes",
        action="store_false",
        help="Do not attach speaker notes.",
    )
    parser.set_defaults(notes=True)

    alt_group = parser.add_mutually_exclusive_group()
    alt_group.add_argument(
        "--alt-text",
        dest="alt_text",
        action="store_true",
        help="Attach image alt text (default).",
    )
    alt_group.add_argument(
        "--no-alt-text",
        dest="alt_text",
        action="store_false",
        help="Do not attach image alt text.",
    )
    parser.set_defaults(alt_text=True)

    return parser.parse_args()


def apply_fade_transition(slide) -> None:
    slide_xml = slide._element
    for child in list(slide_xml):
        if child.tag == qn("p:transition"):
            slide_xml.remove(child)

    transition = OxmlElement("p:transition")
    transition.set("spd", "med")
    transition.append(OxmlElement("p:fade"))

    insert_idx = len(slide_xml)
    for idx, child in enumerate(list(slide_xml)):
        if child.tag in {qn("p:timing"), qn("p:extLst")}:
            insert_idx = idx
            break
    slide_xml.insert(insert_idx, transition)


def set_picture_alt_text(picture, description: str) -> None:
    try:
        c_nv_pr = picture._element.xpath("./p:nvPicPr/p:cNvPr")
        if c_nv_pr:
            c_nv_pr[0].set("descr", description)
    except Exception as err:
        print(f"WARNING: Failed to set alt text: {err}")


def set_speaker_notes(slide, text: str) -> None:
    try:
        slide.notes_slide.notes_text_frame.text = text
    except Exception as err:
        print(f"WARNING: Failed to set speaker notes: {err}")


def convert_with_sips(
    input_path: Path,
    output_path: Path,
    image_format: str,
    jpeg_quality: int,
    max_dimension: int | None,
) -> None:
    if shutil.which("sips") is None:
        raise RuntimeError("sips is required for JPEG conversion or max-dimension scaling on macOS.")
    cmd = ["sips"]
    if max_dimension is not None:
        cmd += ["-Z", str(max_dimension)]
    if image_format == "jpeg":
        cmd += ["-s", "format", "jpeg", "-s", "formatOptions", str(jpeg_quality)]
    else:
        cmd += ["-s", "format", "png"]
    cmd += [str(input_path), "--out", str(output_path)]
    subprocess.run(cmd, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.PIPE, text=True)


def detect_profile_name(args: argparse.Namespace) -> str:
    if args.image_format == "png" and args.max_dimension is None:
        return "master"
    return "delivery"


def default_output_path(out_dir: Path, profile: str) -> Path:
    if profile == "master":
        return out_dir / "Poseidon_AI_MIT_CTO_V3_Visual_First.pptx"
    return out_dir / "Poseidon_AI_MIT_CTO_V3_Visual_First_Delivery.pptx"


def main() -> None:
    args = parse_args()
    profile = detect_profile_name(args)
    jpeg_quality = clamp_jpeg_quality(args.jpeg_quality)
    speaker_notes = load_speaker_notes()
    out_dir = Path(__file__).resolve().parent.parent / "out"
    out_path = args.output or default_output_path(out_dir, profile)

    slide_pngs = [
        "v3-Slide01TitleV3.png",
        slide02_name(args.slide02),
        "v3-Slide03WhyNowV3.png",
        "v3-Slide04SolutionV3.png",
        "v3-Slide05DifferentiationV3.png",
        "v3-Slide06BusinessV3.png",
        "v3-Slide07DemoV3.png",
        "v3-Slide08SummaryV3.png",
        "v3-Slide09EpilogueV3.png",
        "v3-Slide10AppendixV3.png",
        "v3-Slide11FinModelV3.png",
    ]

    missing = [name for name in slide_pngs if not (out_dir / name).exists()]
    if missing:
        print("Missing PNG(s):")
        for name in missing:
            print(f"  - {out_dir / name}")
        raise SystemExit(1)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    needs_processing = args.image_format == "jpeg" or args.max_dimension is not None
    with tempfile.TemporaryDirectory(prefix="poseidon-v3-pptx-") as tmp:
        tmp_dir = Path(tmp)

        for index, name in enumerate(slide_pngs):
            source_png = out_dir / name
            source_image = source_png

            if needs_processing:
                out_ext = ".jpg" if args.image_format == "jpeg" else ".png"
                processed = tmp_dir / f"{source_png.stem}{out_ext}"
                convert_with_sips(
                    input_path=source_png,
                    output_path=processed,
                    image_format=args.image_format,
                    jpeg_quality=jpeg_quality,
                    max_dimension=args.max_dimension,
                )
                source_image = processed

            slide = prs.slides.add_slide(blank)
            picture = slide.shapes.add_picture(
                str(source_image),
                Inches(0),
                Inches(0),
                width=prs.slide_width,
                height=prs.slide_height,
            )

            if args.alt_text:
                set_picture_alt_text(picture, TITLE_FOR.get(name, f"Slide {index + 1}"))

            if args.notes:
                note_text = speaker_notes[index] if index < len(speaker_notes) else ""
                set_speaker_notes(slide, note_text)

            if args.transitions:
                apply_fade_transition(slide)

        out_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out_path))

    size_mb = out_path.stat().st_size / (1024 * 1024)
    print(f"PPTX generated: {out_path} ({size_mb:.1f} MB)")
    print(f"Slides: {len(slide_pngs)}")
    print(
        "Options: "
        f"profile={profile}, "
        f"image_format={args.image_format}, "
        f"jpeg_quality={jpeg_quality}, "
        f"max_dimension={args.max_dimension}, "
        f"notes={args.notes}, "
        f"alt_text={args.alt_text}, "
        f"transitions={args.transitions}"
    )

    latest_marker = out_dir / "LATEST_V3_PPTX.txt"
    latest_marker.write_text(
        (
            f"{out_path.name} | profile={profile} | "
            f"{datetime.now().strftime('%Y%m%d_%H%M%S')}\n"
        ),
        encoding="utf-8",
    )
    print(f"Updated marker: {latest_marker}")


if __name__ == "__main__":
    main()
