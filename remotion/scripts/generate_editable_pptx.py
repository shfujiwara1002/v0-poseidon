#!/usr/bin/env python3
"""
Generate editable PPTX with PNG backgrounds and text-only overlays.
Run from repo root: python scripts/generate_editable_pptx.py
Requires: pip install python-pptx
Input: output/slide_data.json, output/png/*.png
Output: output/poseidon-mit-capstone-final-editable.pptx
"""

import json
from pathlib import Path
from typing import Dict, List, Optional

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
except ImportError:
    print("Run: pip install python-pptx")
    raise SystemExit(1)

REPO_ROOT = Path(__file__).resolve().parent.parent
SLIDE_DATA_JSON = REPO_ROOT / "output" / "slide_data.json"
PNG_DIR = REPO_ROOT / "output" / "png"
OUTPUT_PPTX = REPO_ROOT / "output" / "poseidon-mit-capstone-final-editable.pptx"

# Color palette from theme.ts
COLORS = {
    'cyan': RGBColor(0, 240, 255),  # #00F0FF
    'white': RGBColor(255, 255, 255),
    'lightGray': RGBColor(203, 213, 225),
    'gray': RGBColor(148, 163, 184),
}

# Typography
FONT_NAME = "Space Grotesk"
FONT_FALLBACK = "Arial"

# Slide dimensions (16:9)
SLIDE_WIDTH = 13.333
SLIDE_HEIGHT = 7.5


def add_textbox(slide, left_in, top_in, width_in, height_in, text, font_size_pt, color=None, bold=False, align='left', font_name=None, transparent=True):
    """Add a transparent text box with specified styling."""
    textbox = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))

    # Make textbox transparent
    if transparent:
        textbox.fill.background()
        textbox.line.fill.background()

    text_frame = textbox.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP

    # Paragraph formatting
    paragraph = text_frame.paragraphs[0]
    if align == 'center':
        paragraph.alignment = PP_ALIGN.CENTER
    elif align == 'right':
        paragraph.alignment = PP_ALIGN.RIGHT
    else:
        paragraph.alignment = PP_ALIGN.LEFT

    # Font formatting
    font = paragraph.font
    font.size = Pt(font_size_pt)
    font.name = font_name or FONT_NAME
    font.bold = bold
    if color:
        font.color.rgb = color

    return textbox


def add_background_image(slide, png_path):
    """Add PNG as full-slide background."""
    if png_path.exists():
        left = Inches(0)
        top = Inches(0)
        width = Inches(SLIDE_WIDTH)
        height = Inches(SLIDE_HEIGHT)
        slide.shapes.add_picture(str(png_path), left, top, width=width, height=height)
        return True
    return False


def add_bullet_list(slide, left_in, top_in, width_in, height_in, bullets: List[str], font_size_pt, color=None, font_name=None):
    """Add a transparent bullet list text box."""
    textbox = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))

    # Make transparent
    textbox.fill.background()
    textbox.line.fill.background()

    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP

    for i, bullet in enumerate(bullets):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = bullet
        p.level = 0
        p.alignment = PP_ALIGN.LEFT

        font = p.font
        font.size = Pt(font_size_pt)
        font.name = font_name or FONT_NAME
        if color:
            font.color.rgb = color

    return textbox


def slide_01_title(prs, data, png_path):
    """Slide 01: Title - Text overlays only"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background_image(slide, png_path)

    # Title (center, large)
    add_textbox(slide, 1, 2.5, 11.333, 1, data['title'], 96, COLORS['cyan'], bold=True, align='center')

    # Subtitle
    add_textbox(slide, 1, 3.7, 11.333, 0.6, data['subtitle'], 36, COLORS['white'], align='center')

    # Body
    add_textbox(slide, 1, 4.5, 11.333, 0.5, data['body'], 28, COLORS['lightGray'], align='center')


def slide_02_problem(prs, data, png_path):
    """Slide 02: Problem - Text overlays only"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background_image(slide, png_path)

    # Title (large number)
    add_textbox(slide, 1, 1.2, 11.333, 1, data['title'], 96, COLORS['white'], bold=True, align='center')

    # Subtitle
    add_textbox(slide, 1, 2.4, 11.333, 0.4, data['subtitle'], 28, COLORS['lightGray'], align='center')

    # Body
    add_textbox(slide, 1.5, 3.2, 10.333, 0.8, data['body'], 24, COLORS['white'])


def slide_03_why_now(prs, data, png_path):
    """Slide 03: Why Now - Text overlays only"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background_image(slide, png_path)

    # Title
    add_textbox(slide, 1, 1, 11.333, 0.6, data['title'], 64, COLORS['cyan'], bold=True, align='center')

    # Subtitle
    add_textbox(slide, 1, 1.8, 11.333, 0.4, data['subtitle'], 20, COLORS['lightGray'], align='center')

    # Forces (3 columns) - text only
    if 'forces' in data:
        for i, force in enumerate(data['forces']):
            x = 1.2 + (i * 4)
            y = 3
            w = 3.5

            # Label
            add_textbox(slide, x, y, w, 0.4, force['label'], 22, COLORS['cyan'], bold=True)

            # Bullets
            if 'bullets' in force:
                add_bullet_list(slide, x, y + 0.6, w, 2.5, force['bullets'], 14, COLORS['lightGray'])


def slide_04_solution(prs, data, png_path):
    """Slide 04: Solution - Text overlays only"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background_image(slide, png_path)

    # Title
    add_textbox(slide, 1, 0.8, 11.333, 0.5, data['title'], 56, COLORS['cyan'], bold=True, align='center')

    # Subtitle
    add_textbox(slide, 1, 1.4, 11.333, 0.3, data['subtitle'], 20, COLORS['lightGray'], align='center')

    # Engines (2x2 grid)
    if 'engines' in data:
        positions = [
            (1.2, 2.4), (7.2, 2.4),  # Row 1
            (1.2, 4.7), (7.2, 4.7),  # Row 2
        ]

        for i, engine in enumerate(data['engines'][:4]):
            x, y = positions[i]
            w = 5

            # Name
            add_textbox(slide, x, y, w, 0.4, engine['name'], 28, COLORS['white'], bold=True)

            # Bullets
            if 'bullets' in engine:
                add_bullet_list(slide, x, y + 0.5, w, 1.5, engine['bullets'], 12, COLORS['lightGray'])


def slide_05_trust(prs, data, png_path):
    """Slide 05: Trust - Text overlays only"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background_image(slide, png_path)

    # Title
    add_textbox(slide, 1, 0.8, 11.333, 0.5, data['title'], 56, COLORS['cyan'], bold=True, align='center')

    # Subtitle
    add_textbox(slide, 1, 1.4, 11.333, 0.4, data['subtitle'], 20, COLORS['lightGray'], align='center')

    # Points (2x2 grid)
    if 'points' in data:
        positions = [(1.2, 2.4), (7.2, 2.4), (1.2, 4.7), (7.2, 4.7)]

        for i, point in enumerate(data['points'][:4]):
            x, y = positions[i]
            w = 5

            # Title
            add_textbox(slide, x, y, w, 0.35, point['title'], 18, COLORS['cyan'], bold=True)

            # Description
            add_textbox(slide, x, y + 0.45, w, 1.5, point['desc'], 12, COLORS['lightGray'])


def slide_06_differentiation(prs, data, png_path):
    """Slide 06: Differentiation - Text overlays only"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background_image(slide, png_path)

    # Title
    add_textbox(slide, 1, 0.8, 11.333, 0.5, data['title'], 56, COLORS['cyan'], bold=True, align='center')

    # Subtitle
    add_textbox(slide, 1, 1.4, 11.333, 0.3, data['subtitle'], 18, COLORS['lightGray'], align='center')

    # Narrative
    if 'narrative' in data:
        add_textbox(slide, 1.5, 2, 10.333, 1, data['narrative'], 16, COLORS['white'])


def slide_07_business(prs, data, png_path):
    """Slide 07: Business - Text overlays only"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background_image(slide, png_path)

    # Title
    add_textbox(slide, 1, 0.8, 11.333, 0.5, data['title'], 56, COLORS['cyan'], bold=True, align='center')

    # Subtitle
    add_textbox(slide, 1, 1.4, 11.333, 0.3, data['subtitle'], 18, COLORS['lightGray'], align='center')

    # Metrics
    if 'metrics' in data:
        x_start = 1.5
        y = 2.3
        metric_width = 2

        for i, metric in enumerate(data['metrics']):
            x = x_start + (i * metric_width)
            # Label
            add_textbox(slide, x, y, metric_width, 0.3, metric['label'], 14, COLORS['gray'], align='center')
            # Value
            add_textbox(slide, x, y + 0.35, metric_width, 0.5, metric['value'], 28, COLORS['white'], bold=True, align='center')

    # Revenue streams
    if 'tsstreams' in data:
        for i, stream in enumerate(data['tsstreams']):
            x = 1.5
            y = 3.5 + (i * 1.2)
            w = 10.333

            # Title
            add_textbox(slide, x, y, w, 0.3, stream['title'], 18, COLORS['cyan'], bold=True)
            # Description
            add_textbox(slide, x, y + 0.35, w, 0.7, stream['desc'], 12, COLORS['lightGray'])


def slide_08_roadmap(prs, data, png_path):
    """Slide 08: Roadmap - Text overlays only"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background_image(slide, png_path)

    # Title
    add_textbox(slide, 1, 0.8, 11.333, 0.5, data['title'], 56, COLORS['cyan'], bold=True, align='center')

    # Subtitle
    add_textbox(slide, 1, 1.4, 11.333, 0.3, data['subtitle'], 16, COLORS['lightGray'], align='center')

    # Phases (4 cards)
    if 'phases' in data:
        for i, phase in enumerate(data['phases']):
            x = 0.8 + (i * 3)
            y = 2.5
            w = 2.6

            # Label + Title
            add_textbox(slide, x, y, w, 0.3, f"{phase['label']}: {phase['title']}", 14, COLORS['cyan'], bold=True)

            # Description
            add_textbox(slide, x, y + 0.4, w, 0.8, phase['desc'], 11, COLORS['lightGray'])

            # Bullets
            if 'bullets' in phase:
                add_bullet_list(slide, x, y + 1.3, w, 1.5, phase['bullets'], 9, COLORS['white'])


def slide_09_technology(prs, data, png_path):
    """Slide 09: Technology - Text overlays only"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background_image(slide, png_path)

    # Title
    add_textbox(slide, 1, 0.8, 11.333, 0.5, data['title'], 56, COLORS['cyan'], bold=True, align='center')

    # Subtitle
    add_textbox(slide, 1, 1.4, 11.333, 0.3, data['subtitle'], 18, COLORS['lightGray'], align='center')

    # Details
    if 'details' in data:
        add_textbox(slide, 1.5, 2.1, 10.333, 0.4, data['details'], 16, COLORS['white'])

    # Architecture
    if 'architecture' in data:
        add_textbox(slide, 1.5, 2.7, 10.333, 1.2, data['architecture'], 14, COLORS['lightGray'])


def slide_10_vision(prs, data, png_path):
    """Slide 10: Vision - Text overlays only"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background_image(slide, png_path)

    # Title
    add_textbox(slide, 1, 2, 11.333, 0.6, data['title'], 64, COLORS['cyan'], bold=True, align='center')

    # Subtitle
    add_textbox(slide, 1, 2.8, 11.333, 0.5, data['subtitle'], 22, COLORS['lightGray'], align='center')

    # Quote
    if 'quote' in data:
        add_textbox(slide, 2, 3.8, 9.333, 0.8, data['quote'], 36, COLORS['white'], bold=True, align='center')

    # Supporting
    if 'supporting' in data:
        add_textbox(slide, 2, 5, 9.333, 1, data['supporting'], 18, COLORS['lightGray'], align='center')


def slide_11_demo(prs, data, png_path):
    """Slide 11: Demo - Text overlays only"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background_image(slide, png_path)

    # Title
    add_textbox(slide, 1, 1.5, 11.333, 0.6, data['title'], 56, COLORS['cyan'], bold=True, align='center')

    # Subtitle
    add_textbox(slide, 1, 2.3, 11.333, 0.4, data['subtitle'], 24, COLORS['lightGray'], align='center')

    # Caption
    if 'caption' in data:
        add_textbox(slide, 1, 3, 11.333, 0.5, data['caption'], 28, COLORS['white'], bold=True, align='center')


def slide_12_epilogue(prs, data, png_path):
    """Slide 12: Epilogue - Text overlays only"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background_image(slide, png_path)

    # Title
    add_textbox(slide, 1, 2, 11.333, 0.7, data['title'], 56, COLORS['cyan'], bold=True, align='center')

    # Subtitle
    add_textbox(slide, 1, 2.9, 11.333, 0.4, data['subtitle'], 24, COLORS['lightGray'], align='center')

    # CTA
    if 'cta' in data:
        add_textbox(slide, 1, 4.5, 11.333, 0.6, data['cta']['label'], 32, COLORS['white'], bold=True, align='center')
        if 'url' in data['cta']:
            add_textbox(slide, 1, 5.2, 11.333, 0.3, data['cta']['url'], 18, COLORS['cyan'], align='center')


def main():
    # Load slide data
    if not SLIDE_DATA_JSON.exists():
        print(f"Error: {SLIDE_DATA_JSON} not found.")
        print("Run: node scripts/extract_slide_data.mjs")
        return 1

    with open(SLIDE_DATA_JSON, 'r') as f:
        slide_data = json.load(f)

    print(f"Loaded {len(slide_data['slides'])} slides from {SLIDE_DATA_JSON}")

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH)
    prs.slide_height = Inches(SLIDE_HEIGHT)

    # Slide generators with PNG backgrounds
    slide_generators = [
        slide_01_title,
        slide_02_problem,
        slide_03_why_now,
        slide_04_solution,
        slide_05_trust,
        slide_06_differentiation,
        slide_07_business,
        slide_08_roadmap,
        slide_09_technology,
        slide_10_vision,
        slide_11_demo,
        slide_12_epilogue,
    ]

    png_files = [
        "slide-01-title.png",
        "slide-02-problem.png",
        "slide-03-why-now.png",
        "slide-04-solution.png",
        "slide-05-trust.png",
        "slide-06-differentiation.png",
        "slide-07-business.png",
        "slide-08-roadmap.png",
        "slide-09-technology.png",
        "slide-10-vision.png",
        "slide-11-demo.png",
        "slide-12-epilogue.png",
    ]

    for i, slide_info in enumerate(slide_data['slides']):
        slide_num = slide_info['number']
        slide_name = slide_info['name']
        data = slide_info['data']
        png_path = PNG_DIR / png_files[i]

        print(f"[{slide_num}/12] Generating Slide {slide_num:02d}: {slide_name} (with PNG background)")

        if i < len(slide_generators):
            if png_path.exists():
                slide_generators[i](prs, data, png_path)
            else:
                print(f"  Warning: PNG not found: {png_path}")
        else:
            print(f"  Warning: No generator for slide {slide_num}")

    # Save
    OUTPUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PPTX))
    print(f"\n✅ Editable PPTX saved to: {OUTPUT_PPTX}")
    print(f"   File size: {OUTPUT_PPTX.stat().st_size / (1024*1024):.1f} MB")
    print(f"   Format: PNG backgrounds + editable text overlays")


if __name__ == "__main__":
    main()
