#!/usr/bin/env python3
"""
DEPRECATED: V2 PPTX generator kept for compatibility only.
Prefer: remotion/scripts/gen-v3-pptx.py

Generate V2 PPTX from rendered PNG slides using python-pptx.
This avoids optional PPTX features that can trigger Office repair dialogs.
"""

from __future__ import annotations

import argparse
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    print("Run: pip install python-pptx")
    raise SystemExit(1)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--slide02",
        choices=["optionA", "optionB", "optionC", "default"],
        default="default",
        help="Select Slide02 variant file.",
    )
    return parser.parse_args()


def slide02_name(variant: str) -> str:
    if variant == "optionA":
        return "v2-Slide02ProblemOptionA.png"
    if variant == "optionB":
        return "v2-Slide02ProblemOptionB.png"
    if variant == "optionC":
        return "v2-Slide02ProblemOptionC.png"
    return "v2-Slide02ProblemV2.png"


def main() -> None:
    print("⚠ DEPRECATED: gen-v2-pptx.py is kept for compatibility. Prefer gen-v3-pptx.py.")
    args = parse_args()
    out_dir = Path(__file__).resolve().parent.parent / "out"
    out_path = out_dir / "Poseidon_AI_MIT_CTO_V2_Visual_First.pptx"

    slide_pngs = [
        "v2-Slide01TitleV2.png",
        slide02_name(args.slide02),
        "v2-Slide03WhyNowV2.png",
        "v2-Slide04SolutionV2.png",
        "v2-Slide05DifferentiationV2.png",
        "v2-Slide06BusinessV2.png",
        "v2-Slide07DemoV2.png",
        "v2-Slide08SummaryV2.png",
        "v2-Slide09EpilogueV2.png",
    ]

    missing = [name for name in slide_pngs if not (out_dir / name).exists()]
    if missing:
        print("Missing PNG(s):")
        for name in missing:
            print(f"  - {out_dir / name}")
        raise SystemExit(1)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for name in slide_pngs:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(out_dir / name),
            Inches(0),
            Inches(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    size_mb = out_path.stat().st_size / (1024 * 1024)
    print(f"PPTX generated: {out_path} ({size_mb:.1f} MB)")
    print(f"Slides: {len(slide_pngs)}")


if __name__ == "__main__":
    main()
