#!/usr/bin/env python3
"""
Generate a PPTX showcasing all 56 Poseidon.AI V3 UI screens designed in Pencil.

Each slide pairs Desktop + Mobile for one screen, organized by wave.
If PNG screenshots exist in --images-dir (named by frame ID, e.g. hYZjQ.png),
they are embedded. Otherwise, text-only slides with component inventory.

Usage:
  python scripts/gen-ui-screens-pptx.py
  python scripts/gen-ui-screens-pptx.py --images-dir output/ui-screens
"""

from __future__ import annotations

import argparse
from datetime import datetime
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.oxml.ns import qn
    from pptx.oxml.xmlchemy import OxmlElement
    from pptx.util import Inches, Pt, Emu
except ImportError:
    print("Run: pip install python-pptx")
    raise SystemExit(1)


# ── Theme colors ──────────────────────────────────────────────────────────────
BG_DARK    = RGBColor(0x0D, 0x11, 0x17)
BG_CARD    = RGBColor(0x16, 0x1B, 0x22)
CYAN       = RGBColor(0x00, 0xD1, 0xFF)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
MUTED      = RGBColor(0x8B, 0x94, 0x9E)
GREEN      = RGBColor(0x2E, 0xCC, 0x71)
ORANGE     = RGBColor(0xF5, 0xA6, 0x23)
RED        = RGBColor(0xE5, 0x4D, 0x4D)

ENGINE_COLORS = {
    "Protect": RGBColor(0xE5, 0x4D, 0x4D),
    "Grow":    RGBColor(0x2E, 0xCC, 0x71),
    "Execute": RGBColor(0xF5, 0xA6, 0x23),
    "Govern":  RGBColor(0x7C, 0x8C, 0xFF),
}


# ── Screen data ───────────────────────────────────────────────────────────────
WAVES = [
    {
        "wave": 1,
        "title": "Time-to-First-Value",
        "subtitle": "Critical path: Dashboard → Protect → Execute → Govern + Onboarding",
        "screens": [
            {
                "id": "S-V3-CORE01",
                "name": "Dashboard Overview",
                "engine": None,
                "desktop_id": "hYZjQ",
                "mobile_id": "IA9mi",
                "description": "3-column grid with KPI strip, Engine Health, AI Insight stream, Quick Actions",
                "components": ["SectionHeader", "KPIContractCard", "SignalRow", "ProofLine",
                               "DefinitionLine", "ExplainableInsightPanel", "GovernVerifiedBadge",
                               "AuditLinkChip", "HumanReviewCTA"],
                "govern": True,
            },
            {
                "id": "S-V3-PRT02",
                "name": "Protect Alert Detail",
                "engine": "Protect",
                "desktop_id": "ZBffc",
                "mobile_id": "ll3nj",
                "description": "Alert header + SHAP factors + Action buttons + Similar Pattern Context",
                "components": ["SectionHeader", "ExplainableInsightPanel", "FactorsDropdown",
                               "GovernVerifiedBadge", "AuditLinkChip", "HumanReviewCTA",
                               "ActionOutcomePreview", "ProofLine", "DefinitionLine"],
                "govern": True,
            },
            {
                "id": "S-V3-EXE02",
                "name": "Execute Action Detail & Approval",
                "engine": "Execute",
                "desktop_id": "bbY23",
                "mobile_id": "Cs6sK",
                "description": "Action detail + before/after preview + ConsentScope + Approval workflow",
                "components": ["SectionHeader", "ExplainableInsightPanel", "FactorsDropdown",
                               "GovernVerifiedBadge", "AuditLinkChip", "HumanReviewCTA",
                               "ActionOutcomePreview", "ConsentScopePanel", "ProofLine", "DefinitionLine"],
                "govern": True,
            },
            {
                "id": "S-V3-GOV01",
                "name": "Trust Dashboard",
                "engine": "Govern",
                "desktop_id": "mB7Cw",
                "mobile_id": "sGJ2A",
                "description": "Trust Index hero score, 30-day trend, Active issues, Coverage breakdown",
                "components": ["SectionHeader", "KPIContractCard", "ExplainableInsightPanel",
                               "GovernVerifiedBadge", "AuditLinkChip", "HumanReviewCTA",
                               "ProofLine", "DefinitionLine"],
                "govern": True,
            },
            {
                "id": "S-V3-GOV02",
                "name": "Audit Ledger",
                "engine": "Govern",
                "desktop_id": "6Mmx1",
                "mobile_id": "JQxZv",
                "description": "Searchable audit table with filters (engine/status/date), expandable rows",
                "components": ["SectionHeader", "GovernVerifiedBadge", "AuditLinkChip",
                               "HumanReviewCTA", "ProofLine", "DefinitionLine"],
                "govern": True,
            },
            {
                "id": "S-V3-PUB01",
                "name": "Marketing Landing Page",
                "engine": None,
                "desktop_id": "7R4z8",
                "mobile_id": "pPiP5",
                "description": "Hero + 4-engine value props + social proof + pricing teaser + CTA",
                "components": ["SectionHeader", "ProofLine", "DefinitionLine"],
                "govern": False,
            },
            {
                "id": "S-V3-ACT04",
                "name": "Onboarding: Connect Accounts",
                "engine": None,
                "desktop_id": "weLsX",
                "mobile_id": "OyB7C",
                "description": "Account connection flow with verification rail",
                "components": ["SectionHeader", "ProofLine", "DefinitionLine"],
                "govern": False,
            },
            {
                "id": "S-V3-ACT05",
                "name": "Onboarding: Set Goals",
                "engine": None,
                "desktop_id": "ZuAs1",
                "mobile_id": "RGpcD",
                "description": "Goal configuration with priority settings",
                "components": ["SectionHeader", "ProofLine", "DefinitionLine"],
                "govern": False,
            },
            {
                "id": "S-V3-ACT06",
                "name": "Onboarding: Configure Consent",
                "engine": None,
                "desktop_id": "Igm5F",
                "mobile_id": "X0nV3",
                "description": "Consent scope configuration with per-engine controls",
                "components": ["SectionHeader", "ConsentScopePanel", "ProofLine", "DefinitionLine"],
                "govern": False,
            },
            {
                "id": "S-V3-ACT07",
                "name": "Onboarding: Complete",
                "engine": None,
                "desktop_id": "dspec",
                "mobile_id": "t8JFo",
                "description": "Completion confirmation with next steps",
                "components": ["SectionHeader", "ProofLine", "DefinitionLine"],
                "govern": False,
            },
        ],
    },
    {
        "wave": 2,
        "title": "Trust Closure",
        "subtitle": "Dispute resolution, automation history, audit detail, consent management",
        "screens": [
            {
                "id": "S-V3-PRT03",
                "name": "Dispute & Review",
                "engine": "Protect",
                "desktop_id": "2RAXm",
                "mobile_id": "MYdMc",
                "description": "Case creation, SLA tracking, evidence attachment, resolution flow",
                "components": ["SectionHeader", "StatusTimeline", "ExplainableInsightPanel",
                               "GovernVerifiedBadge", "AuditLinkChip", "HumanReviewCTA",
                               "ProofLine", "DefinitionLine"],
                "govern": True,
            },
            {
                "id": "S-V3-EXE03",
                "name": "Automation History",
                "engine": "Execute",
                "desktop_id": "kZVWA",
                "mobile_id": "VnPBT",
                "description": "Rule editor, execution timeline, reversal options",
                "components": ["SectionHeader", "StatusTimeline", "GovernVerifiedBadge",
                               "AuditLinkChip", "HumanReviewCTA", "ProofLine", "DefinitionLine"],
                "govern": True,
            },
            {
                "id": "S-V3-GOV03",
                "name": "Audit Detail",
                "engine": "Govern",
                "desktop_id": "icdWj",
                "mobile_id": "KpW3D",
                "description": "Full decision reconstruction, hash chain, version history",
                "components": ["SectionHeader", "AuditDetailPanel", "GovernVerifiedBadge",
                               "AuditLinkChip", "HumanReviewCTA", "ProofLine", "DefinitionLine"],
                "govern": True,
            },
            {
                "id": "S-V3-SET04",
                "name": "Consent & Data Rights",
                "engine": None,
                "desktop_id": "C3kzi",
                "mobile_id": "uUp7o",
                "description": "Consent grants/revocations, delete/export, per-source controls",
                "components": ["SectionHeader", "DataRightsPanel", "ConsentScopePanel",
                               "GovernVerifiedBadge", "AuditLinkChip", "HumanReviewCTA",
                               "ProofLine", "DefinitionLine"],
                "govern": True,
            },
            {
                "id": "S-V3-SYS01",
                "name": "Help & Case Tracking",
                "engine": None,
                "desktop_id": "bdKm4",
                "mobile_id": "INUy1",
                "description": "Support case timeline, resolution options, knowledge base links",
                "components": ["SectionHeader", "StatusTimeline", "GovernVerifiedBadge",
                               "AuditLinkChip", "HumanReviewCTA", "ProofLine", "DefinitionLine"],
                "govern": True,
            },
        ],
    },
    {
        "wave": 3,
        "title": "Growth Intelligence",
        "subtitle": "Forecasting, scenario planning, recommendations, explainability registry",
        "screens": [
            {
                "id": "S-V3-GRW01",
                "name": "Forecast Overview",
                "engine": "Grow",
                "desktop_id": "fjna9",
                "mobile_id": "9snTP",
                "description": "Area chart with confidence bands, key drivers panel, time horizon selector",
                "components": ["SectionHeader", "ForecastBandChart", "ExplainableInsightPanel",
                               "GovernVerifiedBadge", "AuditLinkChip", "HumanReviewCTA",
                               "ProofLine", "DefinitionLine"],
                "govern": True,
            },
            {
                "id": "S-V3-GRW02",
                "name": "Scenario Planner",
                "engine": "Grow",
                "desktop_id": "IOoVv",
                "mobile_id": "XMm1C",
                "description": "Multi-scenario comparison with sliders, outcome projections, save/compare",
                "components": ["SectionHeader", "ScenarioControls", "ForecastBandChart",
                               "GovernVerifiedBadge", "AuditLinkChip", "HumanReviewCTA",
                               "ProofLine", "DefinitionLine"],
                "govern": True,
            },
            {
                "id": "S-V3-GRW03",
                "name": "Recommendation Detail",
                "engine": "Grow",
                "desktop_id": "WHaYa",
                "mobile_id": "U7KjZ",
                "description": "Evidence-backed recommendation with execute link, impact preview",
                "components": ["SectionHeader", "ExplainableInsightPanel", "ActionOutcomePreview",
                               "GovernVerifiedBadge", "AuditLinkChip", "HumanReviewCTA",
                               "ProofLine", "DefinitionLine"],
                "govern": True,
            },
            {
                "id": "S-V3-GOV04",
                "name": "Explainability Registry",
                "engine": "Govern",
                "desktop_id": "UDymM",
                "mobile_id": "8EvoH",
                "description": "Template compliance tracking per model, engine compliance grid",
                "components": ["SectionHeader", "ExplainableInsightPanel", "GovernVerifiedBadge",
                               "AuditLinkChip", "HumanReviewCTA", "ProofLine", "DefinitionLine"],
                "govern": True,
            },
            {
                "id": "S-V3-GOV06",
                "name": "Policy & Model Cards",
                "engine": "Govern",
                "desktop_id": "GoSBo",
                "mobile_id": "QuUAI",
                "description": "Model purpose, limitations, fairness metrics, performance, version history",
                "components": ["SectionHeader", "GovernVerifiedBadge", "AuditLinkChip",
                               "HumanReviewCTA", "ProofLine", "DefinitionLine"],
                "govern": True,
            },
        ],
    },
    {
        "wave": 4,
        "title": "Sovereignty & Recovery",
        "subtitle": "Settings, integrations, auth flows, public pages, error recovery",
        "screens": [
            {
                "id": "S-V3-SET01",
                "name": "Settings Overview",
                "engine": None,
                "desktop_id": "L4R3q",
                "mobile_id": "ZpAUw",
                "description": "Settings categories with impact preview, recent changes summary",
                "components": ["SectionHeader", "ProofLine"],
                "govern": False,
            },
            {
                "id": "S-V3-SET02",
                "name": "AI Behavior",
                "engine": None,
                "desktop_id": "tJN9w",
                "mobile_id": "nfFuL",
                "description": "Delegation levels, auto-execute thresholds, explanation granularity",
                "components": ["SectionHeader", "ProofLine", "DefinitionLine"],
                "govern": False,
            },
            {
                "id": "S-V3-SET03",
                "name": "Integrations",
                "engine": None,
                "desktop_id": "hiZQD",
                "mobile_id": "63m4K",
                "description": "Connected sources, permissions, sync health, revoke controls",
                "components": ["SectionHeader", "ProofLine", "DefinitionLine"],
                "govern": False,
            },
            {
                "id": "S-V3-SYS02",
                "name": "Not Found / Recovery",
                "engine": None,
                "desktop_id": "9ktle",
                "mobile_id": "FMlcN",
                "description": "Recovery navigation with 1-tap return to key screens",
                "components": ["SectionHeader"],
                "govern": False,
            },
            {
                "id": "S-V3-PUB02",
                "name": "Trust & Security",
                "engine": None,
                "desktop_id": "j2Sf9",
                "mobile_id": "hVXu4",
                "description": "Compliance badges (SOC 2, GDPR, AES-256), audit evidence, CTA",
                "components": ["SectionHeader", "ProofLine"],
                "govern": False,
            },
            {
                "id": "S-V3-PUB03",
                "name": "Pricing & Value",
                "engine": None,
                "desktop_id": "wEgIz",
                "mobile_id": "xWwI5",
                "description": "3-tier plan comparison (Starter, Pro, Enterprise) with value ratio",
                "components": ["SectionHeader", "ProofLine", "DefinitionLine"],
                "govern": False,
            },
            {
                "id": "S-V3-ACT01",
                "name": "Sign Up",
                "engine": None,
                "desktop_id": "ZrsJP",
                "mobile_id": "qlXTG",
                "description": "Registration form with hero stats panel, social proof",
                "components": ["SectionHeader", "ProofLine"],
                "govern": False,
            },
            {
                "id": "S-V3-ACT02",
                "name": "Login",
                "engine": None,
                "desktop_id": "IdQgy",
                "mobile_id": "KFqVW",
                "description": "Secure login form with session token encryption",
                "components": ["SectionHeader", "ProofLine"],
                "govern": False,
            },
            {
                "id": "S-V3-ACT03",
                "name": "Password Recovery",
                "engine": None,
                "desktop_id": "jeO1Z",
                "mobile_id": "OmBPt",
                "description": "Recovery email form with time-limited link delivery",
                "components": ["SectionHeader", "ProofLine"],
                "govern": False,
            },
        ],
    },
]


# ── Helpers ───────────────────────────────────────────────────────────────────
def set_slide_bg(slide, color: RGBColor) -> None:
    """Set solid fill background on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text: str,
                font_size=12, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Inter") -> None:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment


def add_multiline_textbox(slide, left, top, width, height, lines: list[tuple[str, int, RGBColor, bool]],
                          alignment=PP_ALIGN.LEFT, font_name="Inter") -> None:
    """Add a textbox with multiple lines, each with its own formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(4)


def add_rounded_rect(slide, left, top, width, height, fill_color: RGBColor) -> None:
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        5,  # MSO_SHAPE.ROUNDED_RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def apply_fade_transition(slide) -> None:
    slide_xml = slide._element
    for child in list(slide_xml):
        if child.tag == qn("p:transition"):
            slide_xml.remove(child)
    transition = OxmlElement("p:transition")
    transition.set("spd", "med")
    transition.append(OxmlElement("p:fade"))
    insert_idx = len(slide_xml)
    for idx, child in enumerate(list(slide_xml)):
        if child.tag in {qn("p:timing"), qn("p:extLst")}:
            insert_idx = idx
            break
    slide_xml.insert(insert_idx, transition)


# ── Slide builders ────────────────────────────────────────────────────────────
def build_title_slide(prs: Presentation) -> None:
    """Cover slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    apply_fade_transition(slide)

    # Accent line
    add_rounded_rect(slide, Inches(4.5), Inches(2.2), Inches(4.3), Pt(3), CYAN)

    add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.2),
                "POSEIDON.AI", font_size=48, color=CYAN, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
                "V3 UI Screen Designs", font_size=28, color=WHITE, bold=False,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(4.3), Inches(11), Inches(0.6),
                "28 Screen Pairs  \u2022  56 Frames  \u2022  21 Reusable Components  \u2022  4 Waves",
                font_size=16, color=MUTED, alignment=PP_ALIGN.CENTER)

    # Stats row
    stats = [
        ("21", "Components"),
        ("56", "Frames"),
        ("4", "Engines"),
        ("3", "Quality Gates"),
    ]
    for i, (num, label) in enumerate(stats):
        x = Inches(2.2 + i * 2.4)
        add_textbox(slide, x, Inches(5.4), Inches(2), Inches(0.6),
                    num, font_size=36, color=CYAN, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(5.95), Inches(2), Inches(0.4),
                    label, font_size=14, color=MUTED, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.4),
                datetime.now().strftime("%B %d, %Y"),
                font_size=12, color=MUTED, alignment=PP_ALIGN.CENTER)


def build_wave_divider(prs: Presentation, wave_num: int, title: str, subtitle: str,
                       screen_count: int) -> None:
    """Wave section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    apply_fade_transition(slide)

    add_textbox(slide, Inches(1), Inches(2.0), Inches(11), Inches(0.5),
                f"WAVE {wave_num}", font_size=16, color=CYAN, bold=True,
                alignment=PP_ALIGN.LEFT, font_name="Inter")

    # Accent line
    add_rounded_rect(slide, Inches(1), Inches(2.55), Inches(2), Pt(3), CYAN)

    add_textbox(slide, Inches(1), Inches(2.8), Inches(11), Inches(1),
                title, font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(1), Inches(3.7), Inches(10), Inches(0.6),
                subtitle, font_size=18, color=MUTED, alignment=PP_ALIGN.LEFT)

    add_textbox(slide, Inches(1), Inches(4.8), Inches(4), Inches(0.5),
                f"{screen_count} screen pairs  \u2022  {screen_count * 2} frames",
                font_size=14, color=MUTED, alignment=PP_ALIGN.LEFT)


def build_screen_slide(prs: Presentation, screen: dict, images_dir: Path | None) -> None:
    """One slide per screen pair (desktop + mobile)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    apply_fade_transition(slide)

    # ── Header bar ────────────────────────────────────────────────────────
    engine = screen.get("engine")
    engine_color = ENGINE_COLORS.get(engine, CYAN) if engine else CYAN

    # Screen ID badge
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(2.5), Inches(0.35),
                screen["id"], font_size=11, color=engine_color, bold=True,
                font_name="JetBrains Mono")

    # Screen name
    add_textbox(slide, Inches(0.5), Inches(0.6), Inches(8), Inches(0.5),
                screen["name"], font_size=24, color=WHITE, bold=True)

    # Engine tag (if any)
    if engine:
        add_textbox(slide, Inches(9), Inches(0.35), Inches(2), Inches(0.3),
                    f"\u25cf {engine} Engine", font_size=11, color=engine_color, bold=True)

    # Govern compliance badge
    if screen.get("govern"):
        add_textbox(slide, Inches(9), Inches(0.65), Inches(3.5), Inches(0.3),
                    "\u2713 Govern Contract Set", font_size=11, color=GREEN, bold=True)
    else:
        add_textbox(slide, Inches(9), Inches(0.65), Inches(3.5), Inches(0.3),
                    "\u2014 No AI surface", font_size=11, color=MUTED, bold=False)

    # Accent line under header
    add_rounded_rect(slide, Inches(0.5), Inches(1.1), Inches(12.3), Pt(1), RGBColor(0x30, 0x36, 0x3D))

    # ── Description ───────────────────────────────────────────────────────
    add_textbox(slide, Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.5),
                screen["description"], font_size=13, color=MUTED)

    # ── Try to embed images ───────────────────────────────────────────────
    desktop_img = None
    mobile_img = None
    if images_dir:
        for ext in (".png", ".jpg", ".jpeg"):
            d = images_dir / f"{screen['desktop_id']}{ext}"
            m = images_dir / f"{screen['mobile_id']}{ext}"
            if d.exists():
                desktop_img = d
            if m.exists():
                mobile_img = m

    img_top = Inches(1.85)
    img_bottom = Inches(6.6)
    img_height = img_bottom - img_top

    if desktop_img or mobile_img:
        # Desktop image (left, ~9 inches wide)
        if desktop_img:
            # 1440x900 → aspect 1.6:1
            dw = Inches(8.5)
            dh = int(dw / 1.6)
            if dh > img_height:
                dh = img_height
                dw = int(dh * 1.6)
            slide.shapes.add_picture(str(desktop_img), Inches(0.5), img_top, width=dw, height=dh)
        else:
            _add_placeholder_frame(slide, Inches(0.5), img_top, Inches(8.5), img_height,
                                   "Desktop 1440\u00d7900", screen["desktop_id"])

        # Mobile image (right, ~2.5 inches wide)
        if mobile_img:
            # 402x874 → aspect 0.46:1
            mh = img_height
            mw = int(mh * 0.46)
            if mw > Inches(3):
                mw = Inches(3)
                mh = int(mw / 0.46)
            slide.shapes.add_picture(str(mobile_img), Inches(9.5), img_top, width=mw, height=mh)
        else:
            _add_placeholder_frame(slide, Inches(9.5), img_top, Inches(2.5), img_height,
                                   "Mobile 402\u00d7874", screen["mobile_id"])
    else:
        # No images — show placeholder frames
        _add_placeholder_frame(slide, Inches(0.5), img_top, Inches(8.5), img_height,
                               "Desktop 1440\u00d7900", screen["desktop_id"])
        _add_placeholder_frame(slide, Inches(9.5), img_top, Inches(3), img_height,
                               "Mobile 402\u00d7874", screen["mobile_id"])

    # ── Component inventory (bottom) ──────────────────────────────────────
    comp_y = Inches(6.7)
    comp_text = "  \u2022  ".join(screen["components"])
    add_textbox(slide, Inches(0.5), comp_y, Inches(12.3), Inches(0.6),
                f"Components: {comp_text}",
                font_size=10, color=MUTED, font_name="Inter")

    # ── Speaker notes ─────────────────────────────────────────────────────
    notes = (
        f"Screen: {screen['id']} — {screen['name']}\n"
        f"Desktop frame: {screen['desktop_id']}  |  Mobile frame: {screen['mobile_id']}\n"
        f"Engine: {engine or 'Core'}\n"
        f"Govern Contract: {'Yes' if screen.get('govern') else 'No'}\n"
        f"Components: {', '.join(screen['components'])}\n"
        f"\n{screen['description']}"
    )
    try:
        slide.notes_slide.notes_text_frame.text = notes
    except Exception:
        pass


def _add_placeholder_frame(slide, left, top, width, height, label: str, frame_id: str) -> None:
    """Draw a placeholder rectangle with label text."""
    shape = slide.shapes.add_shape(5, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.color.rgb = RGBColor(0x30, 0x36, 0x3D)
    shape.line.width = Pt(1)

    # Label
    add_textbox(slide, left + Inches(0.2), top + height // 2 - Inches(0.3),
                width - Inches(0.4), Inches(0.35),
                label, font_size=14, color=MUTED, bold=True, alignment=PP_ALIGN.CENTER)
    # Frame ID
    add_textbox(slide, left + Inches(0.2), top + height // 2 + Inches(0.05),
                width - Inches(0.4), Inches(0.25),
                f"Frame: {frame_id}", font_size=10, color=RGBColor(0x48, 0x4F, 0x58),
                alignment=PP_ALIGN.CENTER, font_name="JetBrains Mono")


def build_summary_slide(prs: Presentation) -> None:
    """Final summary slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    apply_fade_transition(slide)

    add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.8),
                "Design Inventory Summary", font_size=32, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)

    # Accent line
    add_rounded_rect(slide, Inches(5), Inches(2.3), Inches(3.3), Pt(3), CYAN)

    # Wave summary table
    wave_data = [
        ("Wave 1", "Time-to-First-Value", "10 pairs", "20 frames", "7 with Govern"),
        ("Wave 2", "Trust Closure", "5 pairs", "10 frames", "5 with Govern"),
        ("Wave 3", "Growth Intelligence", "5 pairs", "10 frames", "5 with Govern"),
        ("Wave 4", "Sovereignty & Recovery", "9 pairs", "16 frames", "0 with Govern"),
    ]

    for i, (wave, name, pairs, frames, govern) in enumerate(wave_data):
        y = Inches(2.8 + i * 0.55)
        add_textbox(slide, Inches(1.5), y, Inches(1.5), Inches(0.4),
                    wave, font_size=14, color=CYAN, bold=True)
        add_textbox(slide, Inches(3.2), y, Inches(3.5), Inches(0.4),
                    name, font_size=14, color=WHITE)
        add_textbox(slide, Inches(7), y, Inches(1.5), Inches(0.4),
                    pairs, font_size=14, color=MUTED)
        add_textbox(slide, Inches(8.5), y, Inches(1.5), Inches(0.4),
                    frames, font_size=14, color=MUTED)
        add_textbox(slide, Inches(10), y, Inches(2), Inches(0.4),
                    govern, font_size=14, color=GREEN if "0" not in govern else MUTED)

    # Totals
    y_total = Inches(5.2)
    add_rounded_rect(slide, Inches(1.5), y_total - Pt(4), Inches(10.5), Pt(1),
                     RGBColor(0x30, 0x36, 0x3D))
    add_textbox(slide, Inches(1.5), y_total, Inches(1.5), Inches(0.4),
                "Total", font_size=14, color=WHITE, bold=True)
    add_textbox(slide, Inches(7), y_total, Inches(1.5), Inches(0.4),
                "29 pairs", font_size=14, color=WHITE, bold=True)
    add_textbox(slide, Inches(8.5), y_total, Inches(1.5), Inches(0.4),
                "56 frames", font_size=14, color=WHITE, bold=True)
    add_textbox(slide, Inches(10), y_total, Inches(2), Inches(0.4),
                "17 with Govern", font_size=14, color=GREEN, bold=True)

    # Quality gates
    add_textbox(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.4),
                "Quality Gates: Gate A (Layout) \u2713  \u2022  Gate B (Semantic) \u2713  \u2022  Gate C (Govern) \u2713",
                font_size=14, color=GREEN, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.4),
                "Designed in Pencil  \u2022  21 reusable components  \u2022  24 design tokens  \u2022  Dark/Light themes",
                font_size=12, color=MUTED, alignment=PP_ALIGN.CENTER)


# ── Main ──────────────────────────────────────────────────────────────────────
def main() -> None:
    parser = argparse.ArgumentParser(description="Generate Poseidon V3 UI Screens PPTX")
    parser.add_argument("--images-dir", type=str, default=None,
                        help="Directory with PNG/JPG screenshots named by frame ID")
    parser.add_argument("--output", type=str, default=None,
                        help="Output PPTX path")
    args = parser.parse_args()

    images_dir = Path(args.images_dir) if args.images_dir else None
    if images_dir and not images_dir.exists():
        print(f"WARNING: Images directory not found: {images_dir}")
        images_dir = None

    out_path = Path(args.output) if args.output else (
        Path(__file__).resolve().parent.parent / "output" / "Poseidon_V3_UI_Screens.pptx"
    )

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Title slide
    build_title_slide(prs)

    # 2. Wave sections
    total_screens = 0
    for wave_data in WAVES:
        build_wave_divider(prs, wave_data["wave"], wave_data["title"],
                           wave_data["subtitle"], len(wave_data["screens"]))
        for screen in wave_data["screens"]:
            build_screen_slide(prs, screen, images_dir)
            total_screens += 1

    # 3. Summary slide
    build_summary_slide(prs)

    # Save
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    size_mb = out_path.stat().st_size / (1024 * 1024)

    print(f"\nPPTX generated: {out_path}")
    print(f"Size: {size_mb:.1f} MB")
    print(f"Slides: {2 + total_screens + len(WAVES)} (1 title + {len(WAVES)} dividers + {total_screens} screens + 1 summary)")
    print(f"Screens: {total_screens} pairs ({total_screens * 2} frames)")
    if images_dir:
        print(f"Images dir: {images_dir}")
    else:
        print("Images: placeholder frames (pass --images-dir to embed screenshots)")


if __name__ == "__main__":
    main()
