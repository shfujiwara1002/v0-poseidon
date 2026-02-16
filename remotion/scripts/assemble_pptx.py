#!/usr/bin/env python3
"""
Assemble deck slide images (PNG or JPEG) into a single PPTX.
Run from repo root: python scripts/assemble_pptx.py [--input output/jpeg]
Requires: pip install python-pptx
Input: output/png/*.png or output/jpeg/*.jpg
Output: output/poseidon-mit-capstone-final.pptx
"""

from __future__ import annotations

import argparse
import json
from pathlib import Path
import shutil
import subprocess
import tempfile
from typing import Optional, Tuple

try:
    from pptx import Presentation
    from pptx.oxml.ns import qn
    from pptx.oxml.xmlchemy import OxmlElement
    from pptx.util import Inches, Pt
except ImportError:
    print("Run: pip install python-pptx")
    raise SystemExit(1)

REPO_ROOT = Path(__file__).resolve().parent.parent
OUTPUT_PPTX = REPO_ROOT / "output" / "poseidon-mit-capstone-final.pptx"
DEBUG_PPTX_LOG = REPO_ROOT / "output" / "debug-pptx.ndjson"


def png_dimensions(path: Path) -> Optional[Tuple[int, int]]:
    """Read PNG width and height from IHDR (bytes 16-23)."""
    try:
        with open(path, "rb") as f:
            buf = f.read(24)
        if len(buf) < 24 or buf[:4] != b"\x89PNG":
            return None
        width = int.from_bytes(buf[16:20], "big")
        height = int.from_bytes(buf[20:24], "big")
        return (width, height)
    except OSError:
        return None

SLIDE_NAMES = [
    "slide-01-title",
    "slide-02-problem",
    "slide-03-why-now",
    "slide-04-solution",
    "slide-05-differentiation",
    "slide-06-business",
    "slide-07-demo",
    "slide-08-summary",
    "slide-09-epilogue",
]


TITLE_FOR = {
    "slide-01-title": "Slide 1 — The Guardian Arrives",
    "slide-02-problem": "Slide 2 — The Coordination Gap",
    "slide-03-why-now": "Slide 3 — Why Now",
    "slide-04-solution": "Slide 4 — The Unified AI Backbone",
    "slide-05-differentiation": "Slide 5 — Beyond Aggregation",
    "slide-06-business": "Slide 6 — Roadmap & Governance",
    "slide-07-demo": "Slide 7 — Introduction Video",
    "slide-08-summary": "Slide 8 — Summary",
    "slide-09-epilogue": "Slide 9 — Epilogue",
}

SPEAKER_NOTES = {
    "slide-01-title": "Poseidon.AI is our trusted AI-native money platform. Today I will walk through the four engines and show how they operate as one system.",
    "slide-02-problem": "The hidden coordination tax is real: around $133 per month in user loss, $12.5B annual fraud loss, and major overdraft burden. Mint proved visibility alone is not enough. Users still act as the integrator.",
    "slide-03-why-now": "Three forces converge now: open banking standards, AI economics improving, and user expectations shifting to proactive guidance. 2025 is the window for AI-native financial services.",
    "slide-04-solution": "Our architecture is a closed loop: Protect, Grow, Execute, Govern. Each engine has a clear role, and together they deliver auditable decisions, not isolated insights.",
    "slide-05-differentiation": "We match baseline fintech features, then differentiate on five critical capabilities, especially governance-by-design. That governance moat is difficult to replicate.",
    "slide-06-business": "Execution is phased across four stages. We track hard gates like 9K users, 70% precision, and 99.9% uptime to prove break-even and operational readiness.",
    "slide-07-demo": "The demo follows a 30-second journey: unified dashboard, Protect with explainability, Execute with human approval, and Govern with audit logging.",
    "slide-08-summary": "The strategy rests on three pillars: governance, architecture, and business model. Economics target about $235 yearly user savings, 77% margin, and roughly 8x value-to-cost.",
    "slide-09-epilogue": "We close with the team and operating principle: deterministic compute, explainable AI, and explicit human approval before action. The QR code links to the live demo.",
}


def clamp_jpeg_quality(raw_quality: int) -> int:
    return max(1, min(100, raw_quality))


def has_sips() -> bool:
    return shutil.which("sips") is not None


def convert_png_to_jpeg(png_path: Path, quality: int, temp_dir: Path) -> Path | None:
    jpg_path = temp_dir / f"{png_path.stem}.jpg"
    cmd = [
        "sips",
        "-s",
        "format",
        "jpeg",
        "-s",
        "formatOptions",
        str(quality),
        str(png_path),
        "--out",
        str(jpg_path),
    ]
    try:
        subprocess.run(cmd, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.PIPE, text=True)
    except (subprocess.CalledProcessError, OSError) as err:
        print(f"⚠ JPEG conversion failed for {png_path.name}: {err}")
        return None
    return jpg_path if jpg_path.exists() else None


def set_picture_alt_text(picture, description: str) -> None:
    try:
        c_nv_pr = picture._element.xpath("./p:nvPicPr/p:cNvPr")
        if c_nv_pr:
            c_nv_pr[0].set("descr", description)
    except Exception as err:
        print(f"⚠ Failed to set alt text: {err}")


def set_speaker_notes(slide, text: str) -> None:
    try:
        slide.notes_slide.notes_text_frame.text = text
    except Exception as err:
        print(f"⚠ Failed to set speaker notes: {err}")


def apply_fade_transition(slide) -> None:
    slide_xml = slide._element
    for child in list(slide_xml):
        if child.tag == qn("p:transition"):
            slide_xml.remove(child)

    transition = OxmlElement("p:transition")
    transition.set("spd", "med")
    transition.append(OxmlElement("p:fade"))

    insert_idx = len(slide_xml)
    for idx, child in enumerate(list(slide_xml)):
        if child.tag in {qn("p:timing"), qn("p:extLst")}:
            insert_idx = idx
            break
    slide_xml.insert(insert_idx, transition)


def main():
    parser = argparse.ArgumentParser(description='Assemble slides into PPTX')
    parser.add_argument('--input', default='output/png', help='Input directory (output/png or output/jpeg)')
    parser.add_argument(
        '--jpeg-quality',
        type=int,
        default=92,
        help='JPEG quality (1-100) when --jpeg conversion is explicitly enabled.',
    )
    parser.add_argument(
        '--transitions',
        action='store_true',
        help='Enable fade transitions (default: off).',
    )

    jpeg_group = parser.add_mutually_exclusive_group()
    jpeg_group.add_argument('--jpeg', dest='jpeg', action='store_true', help='Convert PNG input to JPEG before PPTX embedding.')
    jpeg_group.add_argument('--no-jpeg', dest='jpeg', action='store_false', help='Disable PNG-to-JPEG conversion (default).')
    parser.set_defaults(jpeg=False)

    notes_group = parser.add_mutually_exclusive_group()
    notes_group.add_argument('--notes', dest='notes', action='store_true', help='Attach speaker notes (default).')
    notes_group.add_argument('--no-notes', dest='notes', action='store_false', help='Do not attach speaker notes.')
    parser.set_defaults(notes=True)

    alt_group = parser.add_mutually_exclusive_group()
    alt_group.add_argument('--alt-text', dest='alt_text', action='store_true', help='Attach image alt text (default).')
    alt_group.add_argument('--no-alt-text', dest='alt_text', action='store_false', help='Do not attach image alt text.')
    parser.set_defaults(alt_text=True)

    args = parser.parse_args()

    jpeg_quality = clamp_jpeg_quality(args.jpeg_quality)
    image_dir = REPO_ROOT / args.input
    image_dir.mkdir(parents=True, exist_ok=True)

    # Detect image format (PNG or JPEG)
    is_jpeg = 'jpeg' in args.input.lower() or 'jpg' in args.input.lower()
    extension = '.jpg' if is_jpeg else '.png'
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    slide_w_in = prs.slide_width.inches
    slide_h_in = prs.slide_height.inches

    print(f"\nAssembling PPTX from {image_dir} ({extension} files)...")
    print(f"Options: jpeg={args.jpeg} quality={jpeg_quality}, notes={args.notes}, alt_text={args.alt_text}, transitions={args.transitions}")

    sips_available = has_sips()
    if args.jpeg and extension == '.png' and not sips_available:
        print("⚠ 'sips' not found. PNG files will be embedded directly.")

    with tempfile.TemporaryDirectory(prefix="poseidon-assemble-jpeg-") as tmp:
        tmp_dir = Path(tmp)
        for name in SLIDE_NAMES:
            filename = f"{name}{extension}"
            path = image_dir / filename
            slide = prs.slides.add_slide(blank_layout)

            if path.exists():
                dims = png_dimensions(path) if extension == '.png' else None
                try:
                    DEBUG_PPTX_LOG.parent.mkdir(parents=True, exist_ok=True)
                    with open(DEBUG_PPTX_LOG, "a") as f:
                        f.write(
                            json.dumps(
                                {
                                    "file": filename,
                                    "path": str(path),
                                    "imageWidth": dims[0] if dims else None,
                                    "imageHeight": dims[1] if dims else None,
                                    "slideWidthInches": slide_w_in,
                                    "slideHeightInches": slide_h_in,
                                    "format": "JPEG" if is_jpeg else "PNG"
                                }
                            )
                            + "\n"
                        )
                except OSError:
                    pass

                source_path = path
                if extension == '.png' and args.jpeg and sips_available:
                    maybe_jpg = convert_png_to_jpeg(path, jpeg_quality, tmp_dir)
                    if maybe_jpg is not None:
                        source_path = maybe_jpg
                    else:
                        print(f"  ⚠ JPEG conversion failed, keeping PNG: {filename}")

                left, top = Inches(0), Inches(0)
                picture = slide.shapes.add_picture(str(source_path), left, top, width=prs.slide_width, height=prs.slide_height)
                if args.alt_text:
                    set_picture_alt_text(picture, TITLE_FOR.get(name, name))
                print(f"  ✓ {filename}")
            else:
                print(f"  ⚠️  Missing: {filename} (placeholder created)")
                title = TITLE_FOR.get(name, name)
                tx = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12), Inches(0.8))
                tx.text_frame.text = title
                tx.text_frame.paragraphs[0].font.size = Pt(28)
                tx.text_frame.paragraphs[0].font.bold = True

            if args.notes:
                set_speaker_notes(slide, SPEAKER_NOTES.get(name, ""))
            if args.transitions:
                apply_fade_transition(slide)

    OUTPUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PPTX))
    print("Saved:", OUTPUT_PPTX)


if __name__ == "__main__":
    main()
