#!/usr/bin/env python3
"""
Verify the internal structure of Poseidon_AI_MIT_CTO_V3_Visual_First.pptx.

Inspects the PPTX using python-pptx to confirm:
  - Correct slide count
  - Speaker notes present and non-trivial on every slide
  - Alt text present on every slide image
  - Image format (PNG vs JPEG) per slide

Usage:
  python3 scripts/verify-pptx-contents.py <pptx_path> --expected-slides 11 --json
  python3 scripts/verify-pptx-contents.py <pptx_path>   # human-readable output

Exit code 0 = valid.  Non-zero = structural issue detected.
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.shapes.picture import Picture
except ImportError:
    print("Error: python-pptx not installed.  Run: pip install python-pptx", file=sys.stderr)
    sys.exit(2)


MIN_NOTE_LENGTH = 20  # speaker notes shorter than this are considered trivial


def inspect_pptx(pptx_path: Path, expected_slides: int) -> dict:
    """Return a structured report of PPTX contents."""
    prs = Presentation(str(pptx_path))

    slide_count = len(prs.slides)
    slides_with_notes = 0
    slides_with_alt_text = 0
    png_image_count = 0
    jpeg_image_count = 0
    short_notes = 0
    slides_detail = []

    for idx, slide in enumerate(prs.slides):
        slide_num = idx + 1
        detail = {
            "slideNumber": slide_num,
            "hasNotes": False,
            "noteLength": 0,
            "hasAltText": False,
            "imageFormats": [],
        }

        # Speaker notes
        try:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                detail["hasNotes"] = True
                detail["noteLength"] = len(notes_text)
                slides_with_notes += 1
                if len(notes_text) < MIN_NOTE_LENGTH:
                    short_notes += 1
        except Exception:
            pass

        # Images and alt text
        has_alt_text_for_slide = False
        for shape in slide.shapes:
            if isinstance(shape, Picture):
                # Alt text
                try:
                    c_nv_pr = shape._element.xpath("./p:nvPicPr/p:cNvPr")
                    if c_nv_pr:
                        descr = c_nv_pr[0].get("descr", "")
                        if descr.strip():
                            has_alt_text_for_slide = True
                except Exception:
                    pass

                # Image format detection via content type
                try:
                    content_type = shape.image.content_type
                    if "png" in content_type.lower():
                        detail["imageFormats"].append("PNG")
                        png_image_count += 1
                    elif "jpeg" in content_type.lower() or "jpg" in content_type.lower():
                        detail["imageFormats"].append("JPEG")
                        jpeg_image_count += 1
                    else:
                        detail["imageFormats"].append(content_type)
                except Exception:
                    detail["imageFormats"].append("unknown")

        if has_alt_text_for_slide:
            detail["hasAltText"] = True
            slides_with_alt_text += 1

        slides_detail.append(detail)

    file_size_mb = pptx_path.stat().st_size / (1024 * 1024)

    report = {
        "file": str(pptx_path),
        "fileSizeMB": round(file_size_mb, 1),
        "slideCount": slide_count,
        "expectedSlides": expected_slides,
        "slidesWithNotes": slides_with_notes,
        "slidesWithAltText": slides_with_alt_text,
        "pngImageCount": png_image_count,
        "jpegImageCount": jpeg_image_count,
        "shortNotes": short_notes,
        "slides": slides_detail,
    }

    return report


def print_human_report(report: dict, expected_image_format: str, require_notes: bool) -> bool:
    """Print a human-readable verification report."""
    expected = report["expectedSlides"]
    print(f"\nPPTX Verification Report")
    print(f"{'=' * 50}")
    print(f"File:          {report['file']}")
    print(f"Size:          {report['fileSizeMB']} MB")
    print(f"Slides:        {report['slideCount']} (expected {expected})")
    print(f"Speaker notes: {report['slidesWithNotes']}/{report['slideCount']}")
    print(f"Alt text:      {report['slidesWithAltText']}/{report['slideCount']}")
    print(f"PNG images:    {report['pngImageCount']}")
    print(f"JPEG images:   {report['jpegImageCount']}")
    print()

    failures = []

    if report["slideCount"] != expected:
        failures.append(f"Slide count {report['slideCount']} != expected {expected}")

    if require_notes and report["slidesWithNotes"] != expected:
        missing = []
        for s in report["slides"]:
            if not s["hasNotes"]:
                missing.append(f"Slide {s['slideNumber']}")
        failures.append(f"Missing speaker notes: {', '.join(missing)}")

    if report["slidesWithAltText"] != expected:
        missing = []
        for s in report["slides"]:
            if not s["hasAltText"]:
                missing.append(f"Slide {s['slideNumber']}")
        failures.append(f"Missing alt text: {', '.join(missing)}")

    if expected_image_format == "PNG" and report["jpegImageCount"] > 0:
        failures.append(f"Found {report['jpegImageCount']} JPEG image(s) — expected all PNG")
    if expected_image_format == "JPEG" and report["pngImageCount"] > 0:
        failures.append(f"Found {report['pngImageCount']} PNG image(s) — expected all JPEG")

    if require_notes and report["shortNotes"] > 0:
        short = []
        for s in report["slides"]:
            if s["hasNotes"] and s["noteLength"] < MIN_NOTE_LENGTH:
                short.append(f"Slide {s['slideNumber']} ({s['noteLength']} chars)")
        failures.append(f"Trivially short notes: {', '.join(short)}")

    if failures:
        print("FAILURES:")
        for f in failures:
            print(f"  - {f}")
        print()
    else:
        print("All checks passed.\n")

    # Per-slide detail
    print("Per-slide detail:")
    for s in report["slides"]:
        notes_status = f"notes={s['noteLength']}ch" if s["hasNotes"] else "NO NOTES"
        alt_status = "alt=OK" if s["hasAltText"] else "NO ALT"
        fmt = ",".join(s["imageFormats"]) if s["imageFormats"] else "no-image"
        print(f"  Slide {s['slideNumber']:2d}: {notes_status:20s} {alt_status:10s} {fmt}")
    print()

    return len(failures) == 0


def main() -> None:
    parser = argparse.ArgumentParser(description="Verify PPTX internal structure")
    parser.add_argument("pptx_path", type=Path, help="Path to the PPTX file")
    parser.add_argument("--expected-slides", type=int, default=11, help="Expected number of slides")
    parser.add_argument(
        "--expected-image-format",
        choices=["PNG", "JPEG", "ANY"],
        default="ANY",
        help="Expected image format for embedded slide images.",
    )
    parser.add_argument(
        "--require-notes",
        action="store_true",
        help="Fail when any slide is missing notes or has notes shorter than the minimum length.",
    )
    parser.add_argument(
        "--require-alt-text",
        action="store_true",
        help="Fail when any slide is missing image alt text.",
    )
    parser.add_argument("--json", action="store_true", help="Output JSON instead of human-readable")
    args = parser.parse_args()

    if not args.pptx_path.exists():
        print(f"Error: {args.pptx_path} not found", file=sys.stderr)
        sys.exit(1)

    report = inspect_pptx(args.pptx_path, args.expected_slides)

    if args.json:
        print(json.dumps(report))
    else:
        ok = print_human_report(report, args.expected_image_format, args.require_notes)
        if not ok:
            sys.exit(1)

    failures = []
    expected = args.expected_slides
    if report["slideCount"] != expected:
        failures.append(f"Slide count {report['slideCount']} != expected {expected}")
    if args.require_notes:
        if report["slidesWithNotes"] != expected:
            failures.append(f"Speaker notes missing on {expected - report['slidesWithNotes']} slide(s)")
        if report["shortNotes"] > 0:
            failures.append(f"{report['shortNotes']} speaker note(s) are shorter than {MIN_NOTE_LENGTH} chars")
    if args.require_alt_text and report["slidesWithAltText"] != expected:
        failures.append(f"Alt text missing on {expected - report['slidesWithAltText']} slide(s)")
    if args.expected_image_format == "PNG" and report["jpegImageCount"] > 0:
        failures.append(f"Found JPEG images ({report['jpegImageCount']}) but expected PNG only")
    if args.expected_image_format == "JPEG" and report["pngImageCount"] > 0:
        failures.append(f"Found PNG images ({report['pngImageCount']}) but expected JPEG only")
    if failures:
        for message in failures:
            print(f"Error: {message}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
